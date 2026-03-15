#!/usr/bin/env python3
"""
StowStack Sales Pitch Deck Generator
Creates a 12-slide professional pitch deck using python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Color palette
DARK_BG = RGBColor(2, 6, 23)          # #020617
DARK_BLUE = RGBColor(15, 23, 42)      # #0f172a
CARD_BG = RGBColor(30, 41, 59)        # #1e293b
PRIMARY_ACCENT = RGBColor(16, 185, 129)   # #10b981 (emerald)
SECONDARY_ACCENT = RGBColor(5, 150, 105)  # #059669 (emerald-dark)
WHITE = RGBColor(255, 255, 255)
MUTED_TEXT = RGBColor(148, 163, 184)  # #94a3b8

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

def add_blank_slide():
    """Add a blank slide"""
    blank_layout = prs.slide_layouts[6]  # Blank layout
    return prs.slides.add_slide(blank_layout)

def set_bg_color(slide, color):
    """Set slide background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_text(slide, text, x, y, w, h, size=54, bold=True, color=WHITE, align='left'):
    """Add title text"""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT if align == 'left' else (PP_ALIGN.CENTER if align == 'center' else PP_ALIGN.RIGHT)
    return txBox

def add_body_text(slide, text, x, y, w, h, size=18, color=WHITE, align='left', bold=False, valign=MSO_ANCHOR.TOP):
    """Add body text"""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = PP_ALIGN.LEFT if align == 'left' else (PP_ALIGN.CENTER if align == 'center' else PP_ALIGN.RIGHT)
    return txBox

def add_bullet_list(slide, bullets, x, y, w, h, size=16, color=WHITE):
    """Add bulleted list"""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet_text in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet_text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.level = 0
    return txBox

def add_rectangle(slide, x, y, w, h, fill_color, line_color=None, line_width=None):
    """Add rectangle shape"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

# ============================================================================
# SLIDE 1: TITLE SLIDE
# ============================================================================
slide = add_blank_slide()
set_bg_color(slide, DARK_BG)

# Accent bar on left
add_rectangle(slide, 0, 0, 0.15, 5.625, PRIMARY_ACCENT)

# Main title
title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(1.5))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "StowStack"
p.font.size = Pt(72)
p.font.bold = True
p.font.color.rgb = WHITE

# Tagline 1
add_body_text(slide, "The Demand Engine for Self-Storage", 0.8, 3.1, 8.5, 0.5, size=24, color=PRIMARY_ACCENT, bold=True)

# Tagline 2
add_body_text(slide, "Every ad has its own page.\nEvery move-in traces back to the ad that produced it.", 0.8, 3.7, 8.5, 1, size=16, color=MUTED_TEXT)

# ============================================================================
# SLIDE 2: THE PROBLEM
# ============================================================================
slide = add_blank_slide()
set_bg_color(slide, DARK_BG)

# Title
add_title_text(slide, "The Problem", 0.5, 0.4, 9, 0.7, size=44, color=PRIMARY_ACCENT)

# Subtitle
add_body_text(slide, "You're Spending Money on Ads. You Have No Idea Which Ones Are Filling Units.",
              0.5, 1.1, 9, 0.6, size=20, color=WHITE, bold=True)

# Three pain points
pain_points = [
    ("Your agency reports clicks and impressions.\nNone connect to a signed lease.", 1),
    ("All your traffic goes to the same generic\npage regardless of the ad.", 2.5),
    ("You can't tell which campaign produced\nwhich move-in.", 4)
]

for text, y_pos in pain_points:
    # Colored dot
    add_rectangle(slide, 0.7, y_pos + 0.1, 0.12, 0.12, PRIMARY_ACCENT)
    # Text
    add_body_text(slide, text, 1.1, y_pos, 8.4, 0.8, size=16, color=MUTED_TEXT)

# ============================================================================
# SLIDE 3: THE BROKEN FUNNEL
# ============================================================================
slide = add_blank_slide()
set_bg_color(slide, DARK_BG)

add_title_text(slide, "The Broken Funnel", 0.5, 0.4, 9, 0.7, size=44, color=PRIMARY_ACCENT)

# Left side: Broken funnel
x_left = 0.5
add_body_text(slide, "Without StowStack", x_left, 1.2, 4, 0.4, size=14, color=MUTED_TEXT, bold=True)

steps_broken = ["Ad", "Generic Homepage", "Hope", "???"]
for i, step in enumerate(steps_broken):
    y = 1.7 + i * 0.8
    add_rectangle(slide, x_left + 0.3, y, 3.4, 0.6, CARD_BG)
    add_body_text(slide, step, x_left + 0.5, y + 0.1, 3, 0.4, size=12, color=WHITE, align='center')
    if i < len(steps_broken) - 1:
        add_body_text(slide, "↓", x_left + 1.6, y + 0.65, 1, 0.3, size=20, color=MUTED_TEXT, align='center')

# Right side: StowStack funnel
x_right = 5.2
add_body_text(slide, "With StowStack", x_right, 1.2, 4, 0.4, size=14, color=PRIMARY_ACCENT, bold=True)

steps_fixed = ["Ad", "Custom Landing Page", "Embedded Rental", "Tracked Move-In"]
for i, step in enumerate(steps_fixed):
    y = 1.7 + i * 0.8
    add_rectangle(slide, x_right + 0.3, y, 3.4, 0.6, PRIMARY_ACCENT)
    add_body_text(slide, step, x_right + 0.5, y + 0.1, 3, 0.4, size=12, color=DARK_BG, align='center', bold=True)
    if i < len(steps_fixed) - 1:
        add_body_text(slide, "↓", x_right + 1.6, y + 0.65, 1, 0.3, size=20, color=PRIMARY_ACCENT, align='center')

# ============================================================================
# SLIDE 4: WHAT STOWSTACK IS
# ============================================================================
slide = add_blank_slide()
set_bg_color(slide, DARK_BG)

add_title_text(slide, "What StowStack Is", 0.5, 0.4, 9, 0.7, size=44, color=PRIMARY_ACCENT)
add_body_text(slide, "A closed-loop system from ad click to move-in", 0.5, 1.15, 9, 0.4, size=18, color=MUTED_TEXT, bold=True)

# Four components
components = [
    ("Meta + Google\nPaid Ads", "Demand creation"),
    ("Ad-Specific\nLanding Pages", "Each ad gets its own URL,\nheadline, offer"),
    ("Embedded\nRental Flow", "Reserve/move-in without\nleaving the page"),
    ("Full Attribution", "Which ad → reservation\n→ move-in")
]

x_positions = [0.5, 2.6, 4.7, 6.8]
for i, (title, desc) in enumerate(components):
    x = x_positions[i]
    # Card background
    add_rectangle(slide, x, 2, 1.8, 2.8, CARD_BG)
    # Title
    add_body_text(slide, title, x + 0.15, 2.2, 1.5, 0.8, size=13, color=PRIMARY_ACCENT, bold=True, align='center')
    # Description
    add_body_text(slide, desc, x + 0.15, 3.1, 1.5, 1.3, size=11, color=MUTED_TEXT, align='center')

# ============================================================================
# SLIDE 5: HOW IT WORKS
# ============================================================================
slide = add_blank_slide()
set_bg_color(slide, DARK_BG)

add_title_text(slide, "How It Works", 0.5, 0.4, 9, 0.7, size=44, color=PRIMARY_ACCENT)

# Timeline
timeline_items = [
    ("Day 1", "Facility\nAudit"),
    ("Days 1-2", "Landing\nPages Built"),
    ("Days 2-3", "Campaigns\nLaunch"),
    ("Week 1+", "Move-Ins\nStart")
]

x_start = 0.8
box_width = 1.8
spacing = 0.15

for i, (time, action) in enumerate(timeline_items):
    x = x_start + i * (box_width + spacing)

    # Colored box
    color = PRIMARY_ACCENT if i == 0 else (SECONDARY_ACCENT if i == len(timeline_items) - 1 else CARD_BG)
    add_rectangle(slide, x, 1.8, box_width, 2.2, color if color != CARD_BG else CARD_BG)

    # Time label
    text_color = DARK_BG if color == PRIMARY_ACCENT else WHITE
    add_body_text(slide, time, x + 0.1, 2, box_width - 0.2, 0.4, size=11, color=text_color, bold=True, align='center')

    # Action
    add_body_text(slide, action, x + 0.1, 2.6, box_width - 0.2, 1.2, size=13, color=text_color if color != CARD_BG else PRIMARY_ACCENT, bold=True, align='center')

    # Arrow
    if i < len(timeline_items) - 1:
        add_body_text(slide, "→", x + box_width + spacing * 0.8, 2.9, 0.3, 0.5, size=18, color=PRIMARY_ACCENT, align='center')

# ============================================================================
# SLIDE 6: WHY THIS MATTERS
# ============================================================================
slide = add_blank_slide()
set_bg_color(slide, DARK_BG)

add_title_text(slide, "Why This Matters", 0.5, 0.4, 9, 0.7, size=44, color=PRIMARY_ACCENT)
add_body_text(slide, "The Attribution Loop", 0.5, 1.15, 9, 0.4, size=18, color=MUTED_TEXT, bold=True)

# Key metrics
metrics = [
    ("$38-65", "Cost per Move-In", "(vs $150+ industry avg)"),
    ("15-25%", "Landing Page\nConversion Rate", ""),
    ("3-5x", "ROAS", "(Return on Ad Spend)"),
]

x_positions = [0.7, 3.8, 6.9]
for i, (metric, label, note) in enumerate(metrics):
    x = x_positions[i]

    # Card
    add_rectangle(slide, x, 2, 2.4, 2.2, CARD_BG)

    # Metric (large)
    add_body_text(slide, metric, x + 0.2, 2.3, 2, 0.6, size=28, color=PRIMARY_ACCENT, bold=True, align='center')

    # Label
    add_body_text(slide, label, x + 0.2, 3, 2, 0.6, size=13, color=WHITE, bold=True, align='center')

    # Note
    if note:
        add_body_text(slide, note, x + 0.2, 3.7, 2, 0.5, size=10, color=MUTED_TEXT, align='center')

# Bottom text
add_body_text(slide, "Results from our own facilities", 0.5, 4.8, 9, 0.4, size=12, color=MUTED_TEXT, align='center', bold=True)

# ============================================================================
# SLIDE 7: COMPETITIVE COMPARISON
# ============================================================================
slide = add_blank_slide()
set_bg_color(slide, DARK_BG)

add_title_text(slide, "Competitive Comparison", 0.5, 0.4, 9, 0.7, size=44, color=PRIMARY_ACCENT)

# Table headers
headers = ["Capability", "StowStack", "StorageRankers", "Adverank"]
header_x = [0.5, 3.5, 5.8, 7.8]

for i, header in enumerate(headers):
    add_body_text(slide, header, header_x[i], 1.2, 1.8, 0.4, size=12, color=PRIMARY_ACCENT, bold=True, align='center')

# Capabilities
capabilities = [
    ("Meta Ad Campaigns", "✓", "✓", "✓"),
    ("Google Ads", "✓", "✗", "✗"),
    ("Custom Landing Pages", "✓", "✗", "✗"),
    ("Ad-Specific URLs", "✓", "✗", "✗"),
    ("Embedded Rental Flow", "✓", "✗", "✗"),
    ("Full Attribution", "✓", "✗", "✗"),
    ("Move-In Tracking", "✓", "✗", "✗"),
    ("A/B Testing", "✓", "✗", "✗"),
]

y_start = 1.75
row_height = 0.35

for row_idx, (capability, sw, sr, av) in enumerate(capabilities):
    y = y_start + row_idx * row_height

    # Capability name
    add_body_text(slide, capability, 0.5, y, 2.8, 0.3, size=11, color=MUTED_TEXT)

    # StowStack (green background for ✓)
    if sw == "✓":
        add_rectangle(slide, 3.5, y, 1.8, 0.3, PRIMARY_ACCENT)
        add_body_text(slide, sw, 3.5, y, 1.8, 0.3, size=12, color=DARK_BG, bold=True, align='center')
    else:
        add_body_text(slide, sw, 3.5, y, 1.8, 0.3, size=12, color=MUTED_TEXT, align='center')

    # StorageRankers
    add_body_text(slide, sr, 5.8, y, 1.8, 0.3, size=12, color=MUTED_TEXT, align='center')

    # Adverank
    add_body_text(slide, av, 7.8, y, 1.8, 0.3, size=12, color=MUTED_TEXT, align='center')

# Bottom stats
add_body_text(slide, "StowStack: 13/13  •  StorageRankers: 2/13  •  Adverank: 3/13",
              0.5, 4.95, 9, 0.35, size=11, color=PRIMARY_ACCENT, bold=True, align='center')

# ============================================================================
# SLIDE 8: OPERATOR-FOUNDED
# ============================================================================
slide = add_blank_slide()
set_bg_color(slide, DARK_BG)

# Large quote
quote = slide.shapes.add_textbox(Inches(0.8), Inches(1), Inches(8.4), Inches(2.2))
tf = quote.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '"I built this because I got tired of agencies burning my ad spend on clicks that never turned into move-ins."'
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = PRIMARY_ACCENT
p.alignment = PP_ALIGN.CENTER

# Attribution
add_body_text(slide, "— Blake, Active Self-Storage Operator", 0.8, 3.3, 8.4, 0.5, size=14, color=MUTED_TEXT, align='center')

# Trust points
trust_points = [
    "Active operator with multiple facilities",
    "Tested StowStack on his own units first",
    "Understands NOI, occupancy, seasonal demand"
]

y = 4.1
for point in trust_points:
    add_body_text(slide, "• " + point, 1.5, y, 7, 0.35, size=12, color=MUTED_TEXT)
    y += 0.38

# ============================================================================
# SLIDE 9: PRICING
# ============================================================================
slide = add_blank_slide()
set_bg_color(slide, DARK_BG)

add_title_text(slide, "Pricing", 0.5, 0.4, 9, 0.7, size=44, color=PRIMARY_ACCENT)

# Three pricing tiers
tiers = [
    {
        "name": "Launch",
        "price": "$750/mo",
        "features": [
            "Meta ads",
            "2 landing pages",
            "storEDGE embed",
            "Monthly report"
        ]
    },
    {
        "name": "Growth",
        "price": "$1,500/mo",
        "features": [
            "Meta + Google ads",
            "5 landing pages",
            "Retargeting",
            "A/B testing",
            "Attribution dashboard"
        ]
    },
    {
        "name": "Portfolio",
        "price": "Custom",
        "features": [
            "5+ facilities",
            "Volume discounts",
            "20-35% off",
            "Dedicated support"
        ]
    }
]

x_positions = [0.5, 3.5, 6.5]

for tier_idx, tier in enumerate(tiers):
    x = x_positions[tier_idx]

    # Card background
    is_featured = tier_idx == 1
    card_color = PRIMARY_ACCENT if is_featured else CARD_BG
    add_rectangle(slide, x, 1.2, 2.8, 4, card_color)

    # Tier name
    name_color = DARK_BG if is_featured else PRIMARY_ACCENT
    add_body_text(slide, tier["name"], x + 0.15, 1.4, 2.5, 0.4, size=16, color=name_color, bold=True, align='center')

    # Price
    price_color = DARK_BG if is_featured else PRIMARY_ACCENT
    add_body_text(slide, tier["price"], x + 0.15, 1.9, 2.5, 0.5, size=20, color=price_color, bold=True, align='center')

    # Features
    features_color = DARK_BG if is_featured else MUTED_TEXT
    y = 2.5
    for feature in tier["features"]:
        add_body_text(slide, "• " + feature, x + 0.3, y, 2.3, 0.35, size=10, color=features_color)
        y += 0.38

# ============================================================================
# SLIDE 10: THE ROI MATH
# ============================================================================
slide = add_blank_slide()
set_bg_color(slide, DARK_BG)

add_title_text(slide, "The ROI Math", 0.5, 0.4, 9, 0.7, size=44, color=PRIMARY_ACCENT)

# Calculation box
add_rectangle(slide, 0.5, 1.2, 9, 3.8, CARD_BG)

# Line 1
add_body_text(slide, "Single move-in = $100-150/mo recurring = $1,200-1,800 LTV",
              0.8, 1.5, 8.4, 0.4, size=14, color=WHITE)

# Line 2
add_body_text(slide, "Growth tier at $1,500/mo needs only 5-10 move-ins = 4-12x ROI",
              0.8, 2.0, 8.4, 0.5, size=14, color=PRIMARY_ACCENT, bold=True)

# Line 3
add_body_text(slide, "As data compounds, cost per move-in drops every month",
              0.8, 2.6, 8.4, 0.5, size=14, color=MUTED_TEXT)

# Example boxes
examples = [
    ("Month 1", "0 move-ins\n$1,500 cost", "-$1,500"),
    ("Month 3", "5 move-ins\n$7,500 revenue", "+$6,000"),
    ("Month 6", "12 move-ins\n$18,000 revenue", "+$16,500"),
]

x_positions = [0.8, 3.5, 6.2]
for i, (month, activity, roi) in enumerate(examples):
    x = x_positions[i]
    y = 3.5

    add_body_text(slide, month, x, y, 2.5, 0.3, size=11, color=MUTED_TEXT, bold=True)
    add_body_text(slide, activity, x, y + 0.35, 2.5, 0.6, size=10, color=WHITE)

    roi_color = PRIMARY_ACCENT if "+" in roi else MUTED_TEXT
    add_body_text(slide, roi, x, y + 1, 2.5, 0.35, size=12, color=roi_color, bold=True, align='center')

# ============================================================================
# SLIDE 11: THE VACANCY COST
# ============================================================================
slide = add_blank_slide()
set_bg_color(slide, DARK_BG)

add_title_text(slide, "The Vacancy Cost", 0.5, 0.4, 9, 0.7, size=44, color=PRIMARY_ACCENT)

# Problem statement
add_body_text(slide, "Every empty unit is burning cash.", 0.5, 1.15, 9, 0.4, size=18, color=WHITE, bold=True)

# The math
add_rectangle(slide, 1.5, 1.8, 7, 2.2, CARD_BG)

add_body_text(slide, "15 vacant units × $120/mo × 3 months", 1.8, 2.0, 6.4, 0.4, size=14, color=MUTED_TEXT)
add_body_text(slide, "= $5,400 lost revenue", 1.8, 2.5, 6.4, 0.4, size=16, color=PRIMARY_ACCENT, bold=True)

# Solution
add_rectangle(slide, 1.5, 4.2, 7, 1, PRIMARY_ACCENT)
add_body_text(slide, "StowStack pays for itself with 2-3 filled units", 1.8, 4.35, 6.4, 0.7, size=16, color=DARK_BG, bold=True, align='center')

# ============================================================================
# SLIDE 12: CTA / NEXT STEPS
# ============================================================================
slide = add_blank_slide()
set_bg_color(slide, DARK_BG)

# Accent bar
add_rectangle(slide, 0, 0, 0.15, 5.625, PRIMARY_ACCENT)

# Main CTA
cta_text = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(1))
tf = cta_text.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Get Your Free Facility Audit"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = PRIMARY_ACCENT

# CTA button
add_rectangle(slide, 2.5, 2.4, 5, 0.6, PRIMARY_ACCENT)
add_body_text(slide, "Schedule a 20-minute call", 2.5, 2.45, 5, 0.5, size=16, color=DARK_BG, bold=True, align='center')

# Contact info
add_body_text(slide, "blake@storepawpaw.com  •  (269) 929-8541  •  stowstack.co",
              0.8, 3.4, 8.4, 0.5, size=14, color=MUTED_TEXT, align='center')

# Trust statement
add_body_text(slide, "No credit card required. Let's see if StowStack is right for you.",
              0.8, 4.1, 8.4, 0.6, size=13, color=MUTED_TEXT, align='center')

# ============================================================================
# Save presentation
# ============================================================================
output_path = "/sessions/fervent-great-babbage/mnt/stowstack-app/sales/StowStack_Pitch_Deck.pptx"
prs.save(output_path)
print(f"✓ Presentation saved to: {output_path}")
print(f"✓ Total slides: 12")
